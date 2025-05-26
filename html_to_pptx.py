from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def html_to_pptx(html_file, output_pptx):
    # Read HTML file
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()

    # Parse HTML
    soup = BeautifulSoup(html_content, 'html.parser')

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Get all sections (h1, h2, h3)
    sections = soup.find_all(['h1', 'h2', 'h3'])

    for section in sections:
        # Create a new slide
        slide_layout = prs.slide_layouts[1]  # Using layout with title and content
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = section.get_text().strip()

        # Get content between this section and next section
        content = []
        current = section.next_sibling
        while current and current.name not in ['h1', 'h2', 'h3']:
            if current.name == 'p':
                content.append(current.get_text().strip())
            current = current.next_sibling

        # Add content to slide
        if content:
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.text = '\n'.join(content)

            # Format text
            for paragraph in tf.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(12)

    # Save the presentation
    prs.save(output_pptx)

def main():
    # Get current directory
    current_dir = os.path.dirname(os.path.abspath(__file__))
    
    # Input HTML file
    html_file = os.path.join(current_dir, 'chatbot_training_data.html')
    
    # Output PPTX file
    output_pptx = os.path.join(current_dir, 'chatbot_training_data.pptx')
    
    try:
        html_to_pptx(html_file, output_pptx)
        print(f"Successfully converted {html_file} to {output_pptx}")
    except Exception as e:
        print(f"Error converting file: {str(e)}")

if __name__ == "__main__":
    main() 